import os
import io
import base64
from dotenv import load_dotenv
import streamlit as st
from PIL import Image
from pptx import Presentation
from PyPDF2 import PdfReader
import fitz  # PyMuPDF for PDF image extraction
from langchain.vectorstores import FAISS
from langchain.chains import ConversationChain
from langchain.prompts import ChatPromptTemplate
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_groq import ChatGroq

# Load environment variables
load_dotenv()
groq_api_key = os.getenv("GROQ_API_KEY")
google_api_key = os.getenv("GOOGLE_API_KEY")

# Check for missing API keys
if not google_api_key or not groq_api_key:
    st.error("❌ API keys are missing. Please check your `.env` file.")
    st.stop()

# Load AI models
def load_model(model_name):
    if model_name == "Groq":
        return ChatGroq(groq_api_key=groq_api_key, model_name="llama-3-70b-versatile")
    else:
        return ChatGoogleGenerativeAI(model="gemini-1.5-flash")

# Extract text from PDF
def extract_txt_pdf(file_path):
    reader = PdfReader(file_path)
    text = "".join([page.extract_text() or "" for page in reader.pages])
    return text

# Extract text from PPT
def extract_txt_ppt(file_path):
    ppt = Presentation(file_path)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Extract images from PPT
def extract_images_from_ppt(file_path, output_folder):
    ppt = Presentation(file_path)
    image_paths = []
    for i, slide in enumerate(ppt.slides):
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_filename = os.path.join(output_folder, f"slide_{i}_image_{j}.png")
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes)
                image_paths.append(image_filename)
    return image_paths

# Extract images from PDF
def extract_images_from_pdf(file_path, output_folder):
    doc = fitz.open(file_path)
    image_paths = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_filename = os.path.join(output_folder, f"page_{page_num}_image_{img_index}.png")
            with open(image_filename, "wb") as img_file:
                img_file.write(image_bytes)
            image_paths.append(image_filename)
    return image_paths

# Extract text from images using gemini-1.5-flash
def extract_image_txt(image_path):
    """Extract text from an image using gemini-1.5-flash."""
    model = ChatGoogleGenerativeAI(model="gemini-1.5-flash")

    # Open the image and convert it to base64
    with open(image_path, "rb") as img_file:
        img_base64 = base64.b64encode(img_file.read()).decode("utf-8")

    # Properly formatted input for gemini-1.5-flash
    input_content = [
        {
            "role": "user",
            "content": [
                {"type": "text", "text": "Extract the text from this image."},
                {"type": "image_url", "image_url": f"data:image/png;base64,{img_base64}"}
            ],
        }
    ]

    # Invoke the model and return extracted text
    response = model.invoke(input_content)
    return response.content if response else ""

# Process file (text + images)
def process_file(file_path, output_folder):
    """Process a PPT/PDF file to extract text and images."""
    if file_path.endswith(".pptx"):
        text = extract_txt_ppt(file_path)
        image_paths = extract_images_from_ppt(file_path, output_folder)
    elif file_path.endswith(".pdf"):
        text = extract_txt_pdf(file_path)
        image_paths = extract_images_from_pdf(file_path, output_folder)
    else:
        raise ValueError("Unsupported file format. Use PPTX or PDF.")

    # Extract text from images
    for image_path in image_paths:
        text += "\n" + extract_image_txt(image_path)

    return text

# Initialize FAISS vector store
def initialize_vector_store(text):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    splits = text_splitter.split_text(text)
    return FAISS.from_texts(splits, embeddings)

# Streamlit app
def main():
    st.title("StudNotes")
    vector_store = None  # Initialize here

    # File upload
    uploaded_file = st.sidebar.file_uploader("Upload your notes (PPT or PDF)", type=["pptx", "pdf"])

    if uploaded_file:
        # Save file
        output_folder = "uploads"
        os.makedirs(output_folder, exist_ok=True)
        file_path = os.path.join(output_folder, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        # Process file
        try:
            text = process_file(file_path, output_folder)
            vector_store = initialize_vector_store(text)
            st.success("✅ File processed successfully!")
            st.write("📝 **Extracted Text:**")
            st.text_area("Text Preview", text[:1000], height=200)
        except Exception as e:
            st.error(f"❌ Error processing file: {e}")

    # Query Interface
    st.header("Ask a Question")
    query = st.text_input("Enter your question:")

    if query:
        if vector_store is None:
            st.error("❌ Please upload a file first.")
        else:
            try:
                docs = vector_store.similarity_search(query, k=5)
                st.write("📌 **Top Matches:**")
                for i, doc in enumerate(docs):
                    st.write(f"**Result {i+1}:**")
                    st.write(doc.page_content[:500])  # Limiting text for better readability
                    st.write("---")

                # Generate Response using LLaMA
                llm = load_model("Groq")
                response = llm.invoke(query)
                st.write("💡 **Generated Response:**")
                st.write(response.content)

            except Exception as e:
                st.error(f"❌ Error querying notes: {e}")

if __name__ == "__main__":
    main()